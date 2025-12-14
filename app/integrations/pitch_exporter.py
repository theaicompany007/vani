"""Pitch deck export functionality (PDF and PowerPoint)"""
import logging
import io
from typing import Dict, Any, Optional
from datetime import datetime

logger = logging.getLogger(__name__)

# Try weasyprint first (requires GTK3 on Windows)
try:
    from weasyprint import HTML, CSS
    WEASYPRINT_AVAILABLE = True
except (ImportError, OSError) as e:
    WEASYPRINT_AVAILABLE = False
    logger.warning(f"weasyprint not available: {e}. PDF export will use reportlab fallback.")

# Try reportlab as fallback (Windows-friendly)
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False
    logger.warning("reportlab not available. PDF export will not work.")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. PowerPoint export will not work.")


class PitchExporter:
    """Export pitch decks to PDF and PowerPoint formats"""
    
    @staticmethod
    def export_to_pdf(pitch_content: Dict[str, Any], company_name: str, contact_name: Optional[str] = None) -> bytes:
        """
        Export pitch to PDF format
        
        Args:
            pitch_content: Dictionary with pitch content (title, problem, solution, hit_list, trojan_horse)
            company_name: Company name
            contact_name: Optional contact name
            
        Returns:
            PDF bytes
        """
        # Try weasyprint first
        if WEASYPRINT_AVAILABLE:
            try:
                html_content = PitchExporter._generate_pdf_html(pitch_content, company_name, contact_name)
                pdf_bytes = HTML(string=html_content).write_pdf()
                logger.info(f"Successfully generated PDF using weasyprint for {company_name}")
                return pdf_bytes
            except Exception as e:
                logger.warning(f"weasyprint failed: {e}, trying alternative method")
        
        # Fallback to reportlab
        if REPORTLAB_AVAILABLE:
            try:
                return PitchExporter._export_to_pdf_reportlab(pitch_content, company_name, contact_name)
            except Exception as e:
                logger.error(f"reportlab PDF generation failed: {e}")
                raise
        
        raise ImportError(
            "PDF export not available. Install weasyprint (requires GTK3 on Windows) or reportlab. "
            "For Windows, consider: pip install reportlab"
        )
    
    @staticmethod
    def _export_to_pdf_reportlab(pitch_content: Dict[str, Any], company_name: str, contact_name: Optional[str] = None) -> bytes:
        """Export to PDF using reportlab (Windows-friendly alternative)"""
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib.colors import HexColor
        
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=1*inch, bottomMargin=1*inch)
        
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=HexColor('#4F46E5'),
            spaceAfter=30,
            alignment=1  # Center
        )
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=18,
            textColor=HexColor('#4F46E5'),
            spaceAfter=12,
            spaceBefore=20
        )
        
        story = []
        
        # Title
        title = pitch_content.get('title', f'Strategic Pitch for {company_name}')
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 0.3*inch))
        
        # Subtitle
        date_str = datetime.now().strftime("%B %d, %Y")
        subtitle = f"{company_name}<br/>{contact_name if contact_name else 'Executive Team'}<br/>{date_str}"
        story.append(Paragraph(subtitle, styles['Normal']))
        story.append(Spacer(1, 0.5*inch))
        
        # Problem
        story.append(Paragraph("The Challenge", heading_style))
        problem = pitch_content.get('problem', '').replace('\n', '<br/>')
        story.append(Paragraph(problem, styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        # Solution
        story.append(Paragraph("The Solution: VANI", heading_style))
        solution = pitch_content.get('solution', '').replace('\n', '<br/>')
        story.append(Paragraph(solution, styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        # Value Proposition
        story.append(Paragraph(f"Value Proposition for {company_name}", heading_style))
        hit_list = pitch_content.get('hit_list', '').replace('\n', '<br/>')
        story.append(Paragraph(hit_list, styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        # Trojan Horse
        story.append(Paragraph("Low-Risk Pilot Approach", heading_style))
        trojan_horse = pitch_content.get('trojan_horse', '').replace('\n', '<br/>')
        story.append(Paragraph(trojan_horse, styles['Normal']))
        
        doc.build(story)
        buffer.seek(0)
        
        logger.info(f"Successfully generated PDF using reportlab for {company_name}")
        return buffer.read()
    
    @staticmethod
    def export_to_pptx(pitch_content: Dict[str, Any], company_name: str, contact_name: Optional[str] = None) -> bytes:
        """
        Export pitch to PowerPoint format
        
        Args:
            pitch_content: Dictionary with pitch content
            company_name: Company name
            contact_name: Optional contact name
            
        Returns:
            PowerPoint bytes
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PowerPoint export. Install with: pip install python-pptx")
        
        try:
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Slide 1: Title Slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            
            title.text = pitch_content.get('title', f'Strategic Pitch for {company_name}')
            subtitle.text = f"{company_name}\n{contact_name if contact_name else 'Executive Team'}\n{datetime.now().strftime('%B %d, %Y')}"
            
            # Slide 2: Problem Statement
            problem_slide = prs.slides.add_slide(prs.slide_layouts[1])
            problem_title = problem_slide.shapes.title
            problem_content = problem_slide.placeholders[1]
            
            problem_title.text = "The Challenge"
            problem_content.text = pitch_content.get('problem', 'Industry-specific challenges')
            
            # Slide 3: Solution Overview
            solution_slide = prs.slides.add_slide(prs.slide_layouts[1])
            solution_title = solution_slide.shapes.title
            solution_content = solution_slide.placeholders[1]
            
            solution_title.text = "The Solution: VANI"
            solution_content.text = pitch_content.get('solution', 'AI-powered solution')
            
            # Slide 4: Hit List / Value Proposition
            hitlist_slide = prs.slides.add_slide(prs.slide_layouts[1])
            hitlist_title = hitlist_slide.shapes.title
            hitlist_content = hitlist_slide.placeholders[1]
            
            hitlist_title.text = f"Value Proposition for {company_name}"
            hitlist_content.text = pitch_content.get('hit_list', 'Company-specific value proposition')
            
            # Slide 5: Trojan Horse Strategy
            trojan_slide = prs.slides.add_slide(prs.slide_layouts[1])
            trojan_title = trojan_slide.shapes.title
            trojan_content = trojan_slide.placeholders[1]
            
            trojan_title.text = "Low-Risk Pilot Approach"
            trojan_content.text = pitch_content.get('trojan_horse', 'Pilot program details')
            
            # Slide 6: Next Steps
            nextsteps_slide = prs.slides.add_slide(prs.slide_layouts[1])
            nextsteps_title = nextsteps_slide.shapes.title
            nextsteps_content = nextsteps_slide.placeholders[1]
            
            nextsteps_title.text = "Next Steps"
            nextsteps_content.text = "• Schedule a discovery call\n• Define pilot scope\n• Set success metrics\n• Launch pilot program"
            
            # Save to bytes
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            
            logger.info(f"Successfully generated PowerPoint for {company_name}")
            return output.read()
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint: {e}")
            raise
    
    @staticmethod
    def _generate_pdf_html(pitch_content: Dict[str, Any], company_name: str, contact_name: Optional[str] = None) -> str:
        """Generate HTML content for PDF export"""
        title = pitch_content.get('title', f'Strategic Pitch for {company_name}')
        problem = pitch_content.get('problem', '')
        solution = pitch_content.get('solution', '')
        hit_list = pitch_content.get('hit_list', '')
        trojan_horse = pitch_content.get('trojan_horse', '')
        
        date_str = datetime.now().strftime("%B %d, %Y")
        
        html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <style>
                @page {{
                    size: A4;
                    margin: 1in;
                }}
                body {{
                    font-family: 'Arial', sans-serif;
                    line-height: 1.6;
                    color: #333;
                }}
                .header {{
                    text-align: center;
                    margin-bottom: 2em;
                    border-bottom: 3px solid #4F46E5;
                    padding-bottom: 1em;
                }}
                .header h1 {{
                    color: #4F46E5;
                    font-size: 28px;
                    margin: 0;
                }}
                .header .subtitle {{
                    color: #666;
                    font-size: 14px;
                    margin-top: 0.5em;
                }}
                .section {{
                    margin-bottom: 2em;
                    page-break-inside: avoid;
                }}
                .section h2 {{
                    color: #4F46E5;
                    font-size: 20px;
                    border-bottom: 2px solid #E5E7EB;
                    padding-bottom: 0.5em;
                    margin-bottom: 1em;
                }}
                .section p {{
                    font-size: 14px;
                    text-align: justify;
                    margin-bottom: 1em;
                }}
                .footer {{
                    margin-top: 3em;
                    text-align: center;
                    font-size: 12px;
                    color: #666;
                    border-top: 1px solid #E5E7EB;
                    padding-top: 1em;
                }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1>{title}</h1>
                <div class="subtitle">
                    {company_name}<br>
                    {contact_name if contact_name else 'Executive Team'}<br>
                    {date_str}
                </div>
            </div>
            
            <div class="section">
                <h2>The Challenge</h2>
                <p>{problem.replace(chr(10), '<br>')}</p>
            </div>
            
            <div class="section">
                <h2>The Solution: VANI</h2>
                <p>{solution.replace(chr(10), '<br>')}</p>
            </div>
            
            <div class="section">
                <h2>Value Proposition for {company_name}</h2>
                <p>{hit_list.replace(chr(10), '<br>')}</p>
            </div>
            
            <div class="section">
                <h2>Low-Risk Pilot Approach</h2>
                <p>{trojan_horse.replace(chr(10), '<br>')}</p>
            </div>
            
            <div class="footer">
                <p>Prepared by Project VANI | Virtual Agent Network Interface</p>
                <p>For more information, contact: vani@theaicompany.ngrok.app</p>
            </div>
        </body>
        </html>
        """
        
        return html

